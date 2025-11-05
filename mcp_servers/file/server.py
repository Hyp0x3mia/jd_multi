"""
PDF / Image / PPTX + LLM 内容分析 MCP 服务（在原 doc_analyze 基础上增强）
- 保持单一工具：doc_analyze（不新增别的工具）
- 新增支持：.pptx 直读文本（文本框/表格），带长文档分片
- 仍支持：pdf（渲染为图片交给 VLM）、jpg/png/bmp（直接喂 VLM）

依赖：
pip install mcp-server-fastmcp pdfplumber pillow python-pptx pymupdf openai python-dotenv
"""

import os
import io
import time
import random
from typing import Optional, List, Tuple
from mcp.server.fastmcp import FastMCP
from pydantic import Field
from pathlib import Path
import base64
from openai import OpenAI
import fitz  # PyMuPDF
from PIL import Image
import re
from dotenv import load_dotenv
import logging, time
log = logging.getLogger(__name__)

try:
    from pptx import Presentation
except Exception:
    Presentation = None  # 若未安装 python-pptx，会在运行时报错提示

# ======================= VLM 指令 =======================
_DEEP_ANALYZER_INSTRUCTION = """You should step-by-step analyze the task and/or the attached content.
* When the task involves playing a game or performing calculations. Please consider the conditions imposed by the game or calculation rules. You may take extreme conditions into account.
* When the task involves spelling words, you must ensure that the spelling rules are followed and that the resulting word is meaningful.
* When the task involves compute the area in a specific polygon. You should separate the polygon into sub-polygons and ensure that the area of each sub-polygon is computable (e.g, rectangle, circle, triangle, etc.). Step-by-step to compute the area of each sub-polygon and sum them up to get the final area.
* When the task involves calculation and statistics, it is essential to consider all constraints. Failing to account for these constraints can easily lead to statistical errors.

Here is the task:
"""

# ======================= 初始化 =======================
mcp = FastMCP()
load_dotenv("./.env")

# 安全读取环境变量，避免 None.rstrip()
_VL_URL = os.getenv("VL_URL") 
_VL_KEY = os.getenv("VL_KEY")
_VL_MODEL = os.getenv("VL_MODEL")
_DOC_MAX_WIDTH = int(os.getenv("DOC_MAX_WIDTH", "1600"))
_PDF_RENDER_DPI = int(os.getenv("PDF_RENDER_DPI", "220"))
_MAX_CHARS_PER_CALL = int(os.getenv("MAX_CHARS_PER_CALL", "7000"))  # 文本分片上限，控制单次 tokens

if not _VL_URL or not _VL_KEY or not _VL_MODEL:
    raise RuntimeError("缺少必需环境变量：VL_URL / VL_KEY / VL_MODEL")

_VL_URL = _VL_URL.rstrip("/")
client = OpenAI(base_url=_VL_URL, api_key=_VL_KEY)
model_name = _VL_MODEL

# ======================= 通用工具 =======================
def _resize_image_bytes(img_bytes: bytes, max_width: int = _DOC_MAX_WIDTH) -> bytes:
    """将图片按最长边等比缩放到不超过 max_width，返回PNG字节。"""
    try:
        with Image.open(io.BytesIO(img_bytes)) as im:
            im = im.convert("RGB")
            w, h = im.size
            if w <= max_width:
                buf = io.BytesIO()
                im.save(buf, format="PNG", optimize=True)
                return buf.getvalue()
            scale = max_width / float(w)
            new_size = (max_width, int(h * scale))
            im = im.resize(new_size, Image.LANCZOS)
            buf = io.BytesIO()
            im.save(buf, format="PNG", optimize=True)
            return buf.getvalue()
    except Exception:
        return img_bytes

def _call_vlm(messages: List[dict], max_retries: int = 4, base_sleep: float = 0.6, timeout: int = 600) -> str:
    """稳定调用 VLM，带指数退避重试。"""
    last_err = None
    for attempt in range(1, max_retries + 1):
        try:
            resp = client.chat.completions.create(
                model=model_name,
                messages=messages,
                temperature=0,
                top_p=0.8,
                max_tokens=512,
                stream=False,
                timeout=timeout,
            )
            return resp.choices[0].message.content
        except Exception as e:
            last_err = e
            sleep_s = base_sleep * (1.6 ** (attempt - 1)) + random.uniform(0, 0.2)
            time.sleep(sleep_s)
    raise RuntimeError(f"VLM call failed after {max_retries} attempts: {last_err}")

# ======================= PDF / Image 走 VLM OCR =======================
def _analyze_image(path: str, prompt: str) -> str:
    try:
        with open(path, "rb") as f:
            raw = f.read()
        img_bytes = _resize_image_bytes(raw, max_width=_DOC_MAX_WIDTH)
        b64_image = base64.b64encode(img_bytes).decode("utf-8")

        content = [
            {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64_image}"}},
            {"type": "text", "text": (_DEEP_ANALYZER_INSTRUCTION + (prompt or "请只输出答案，禁止解释。"))},
        ]
        messages = [{"role": "user", "content": content}]
        return _call_vlm(messages)
    except Exception as e:
        return f"[IMAGE analyze failed: {e}]"

def _analyze_pdf(path: str, prompt: str) -> str:
    try:
        doc = fitz.open(path)
        content: List[dict] = []
        for page in doc:
            pix = page.get_pixmap(dpi=_PDF_RENDER_DPI, alpha=False)
            png_bytes = pix.tobytes("png")
            png_bytes = _resize_image_bytes(png_bytes, max_width=_DOC_MAX_WIDTH)
            b64_image = base64.b64encode(png_bytes).decode()
            content.append({"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64_image}"}})
        doc.close()

        content.append({"type": "text", "text": (_DEEP_ANALYZER_INSTRUCTION + (prompt or "请只输出答案，禁止解释。"))})
        messages = [{"role": "user", "content": content}]
        return _call_vlm(messages)
    except Exception as e:
        return f"[PDF analyze failed: {e}]"

# ======================= PPTX 直读文本 =======================
def _extract_pptx_text(path: str) -> List[Tuple[int, str]]:
    """
    返回 [(slide_no, slide_text)]，只读文本框与表格；不做图片 OCR（保持与“无新增工具”一致）。
    """
    if Presentation is None:
        raise RuntimeError("缺少依赖 python-pptx：pip install python-pptx")

    prs = Presentation(path)
    results: List[Tuple[int, str]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        chunks: List[str] = []
        for shape in slide.shapes:
            # 文本框 / 自带 text_frame 的形状
            if hasattr(shape, "has_text_frame") and shape.has_text_frame and shape.text_frame:
                paras = [p.text for p in shape.text_frame.paragraphs if p.text]
                if paras:
                    chunks.append("\n".join(paras))
            # 表格
            if hasattr(shape, "has_table") and shape.has_table and shape.table:
                for r in shape.table.rows:
                    for c in r.cells:
                        if c.text:
                            chunks.append(c.text)
        text = "\n".join([t for t in chunks if t]).strip()
        results.append((idx, text))
    return results

def _pptx_find_keyword_positions(slide_texts: List[Tuple[int, str]], keyword: str) -> List[Tuple[int, int]]:
    """
    返回 [(slide_no, index_in_text)]，表示关键字首次出现的位置。
    index_in_text 为文本中的字符索引，-1 表示未出现（不会返回）。
    """
    positions: List[Tuple[int, int]] = []
    if not keyword:
        return positions
    for slide_no, txt in slide_texts:
        idx = txt.find(keyword)
        if idx >= 0:
            positions.append((slide_no, idx))
    return positions

def _pptx_get_following_lines(slide_texts: List[Tuple[int, str]], slide_no: int, start_index: int, max_lines: int = 80) -> str:
    """
    从指定页的 start_index 开始，提取该页后续行，必要时继续拼接下一页，最多返回 max_lines 行。
    行按原文本换行切分。
    """
    # 建立页码索引
    page_map = {no: txt for no, txt in slide_texts}
    lines_out: List[str] = []
    cur_no = slide_no
    cur_idx = start_index
    while len(lines_out) < max_lines and cur_no in page_map:
        txt = page_map[cur_no]
        seg = txt[cur_idx:]
        seg_lines = [ln for ln in seg.splitlines() if ln is not None]
        for ln in seg_lines:
            if len(lines_out) >= max_lines:
                break
            lines_out.append(ln)
        cur_no += 1
        cur_idx = 0
    return "\n".join(lines_out)

def _chunk_texts(slide_texts: List[Tuple[int, str]], max_chars: int = _MAX_CHARS_PER_CALL) -> List[str]:
    """
    将全量 PPT 文本按字数拼接成多个分片。
    每个分片使用格式：
    【第n页】
    <文本>
    ----
    """
    chunks: List[str] = []
    cur: List[str] = []
    cur_len = 0
    for slide_no, txt in slide_texts:
        block = f"【第{slide_no}页】\n{txt}\n----\n"
        blen = len(block)
        if cur_len + blen > max_chars and cur:
            chunks.append("".join(cur))
            cur = [block]
            cur_len = blen
        else:
            cur.append(block)
            cur_len += blen
    if cur:
        chunks.append("".join(cur))
    return chunks

def _analyze_pptx(path: str, prompt: str) -> str:
    """
    读取 PPTX 文本 → 长文档自动分片 → 多次调用 VLM → 合并答案为最终输出
    """
    slides = _extract_pptx_text(path)  # [(no, text)]
    # 给 VLM 看的内容，显式标记页码，便于它做“第几页/相同页码匹配/百分比搜索”等操作
    parts = _chunk_texts(slides, max_chars=_MAX_CHARS_PER_CALL)

    final_notes: List[str] = []
    # 如果只有一个分片，直接一次调用
    if len(parts) == 1:
        messages = [{"role": "user", "content": [
            {"type": "text", "text": f"{_DEEP_ANALYZER_INSTRUCTION}{prompt or '请只输出答案，禁止解释。'}\n\n<<<PPT-CONTENT>>>\n{parts[0]}"}
        ]}]
        return _call_vlm(messages)

    # 多分片：对每个分片单独问同一个任务，再把中间答案交给最后一轮归并
    for i, seg in enumerate(parts, start=1):
        sub_prompt = (prompt or "请只输出答案，禁止解释。") + f"\n(注：这是第 {i}/{len(parts)} 个内容分片)"
        messages = [{"role": "user", "content": [
            {"type": "text", "text": f"{_DEEP_ANALYZER_INSTRUCTION}{sub_prompt}\n\n<<<PPT-SEGMENT {i}>>>\n{seg}"}
        ]}]
        ans_i = _call_vlm(messages)
        final_notes.append(f"[SEG-{i}] {ans_i}")

    # 归并阶段：把每个分片答案提交给模型，让它综合成单一最终答案（仍保持“只输出答案”的习惯）
    merge_messages = [{"role": "user", "content": [
        {"type": "text", "text": (
            f"{_DEEP_ANALYZER_INSTRUCTION}"
            f"你将收到多个内容分片的中间答案，请综合它们，**只输出最终答案**，不要解释：\n\n"
            + "\n".join(final_notes)
        )}
    ]}]
    return _call_vlm(merge_messages)

# ======================= MCP 工具 =======================
@mcp.tool(description="读取本地 PDF / 图片 / PPTX，提取内容并调用 VLM 做分析（通用入口）")
async def doc_analyze(
    path: str = Field(description="本地文件路径（pdf / pptx / jpg / jpeg / png / bmp）"),
    prompt: Optional[str] = Field(
        default=None, description="可选：给 LLM 的额外提示，如“输出第一张不含英文的页码”“查找最大百分比并只输出数值”等"
    ),
):
    log.info("[doc_analyze] tool-timeout=300, vlm-timeout=120")
    # 兼容相对路径 → 工作根路径
    path_str = str(path).strip()
    if path_str:
        if path_str.startswith('./') or (not os.path.isabs(path_str) and '/' not in path_str):
            here = Path(__file__).resolve()
            root_path = here.parent.parent.parent
            if root_path.exists():
                path = str(root_path / path_str.lstrip('./'))

    if not os.path.exists(path):
        return {"error": f"文件不存在: {path}"}

    ext = os.path.splitext(path)[-1].lower()

    try:
        if ext == ".pdf":
            analysis = _analyze_pdf(path, prompt)
        elif ext in {".jpg", ".jpeg", ".png", ".bmp"}:
            analysis = _analyze_image(path, prompt)
        elif ext == ".pptx":
            analysis = _analyze_pptx(path, prompt)
        else:
            return {"error": "仅支持 pdf / pptx / jpg / jpeg / png / bmp"}
    except Exception as e:
        return {"error": "分析失败", "raw": f"{type(e).__name__}: {e}"}

    if not analysis or (analysis.startswith("[") and "failed" in analysis):
        return {"error": "分析失败", "raw": analysis}

    return {
        "status": "success",
        "path": path,
        "analysis": analysis,
    }

@mcp.tool(description="提取 PPTX 所有页的纯文本（每页一段），供上层步骤化处理")
async def pptx_extract_text(path: str = Field(description="本地 PPTX 文件路径")):
    if not os.path.exists(path):
        return {"error": f"文件不存在: {path}"}
    if os.path.splitext(path)[-1].lower() != ".pptx":
        return {"error": "仅支持 PPTX"}
    try:
        slides = _extract_pptx_text(path)
        return {"status": "success", "slides": [{"slide": no, "text": txt} for no, txt in slides]}
    except Exception as e:
        return {"error": "提取失败", "raw": f"{type(e).__name__}: {e}"}

@mcp.tool(description="在 PPTX 中查找关键字出现的页与上下文摘要")
async def pptx_find_keyword(
    path: str = Field(description="本地 PPTX 文件路径"),
    keyword: str = Field(description="要查找的关键字，如：文艺表演节目单"),
    context_lines: int = Field(default=40, description="返回关键字处开始的后续若干行，便于解析列表")
):
    if not os.path.exists(path):
        return {"error": f"文件不存在: {path}"}
    if os.path.splitext(path)[-1].lower() != ".pptx":
        return {"error": "仅支持 PPTX"}
    try:
        slides = _extract_pptx_text(path)
        poses = _pptx_find_keyword_positions(slides, keyword)
        if not poses:
            return {"status": "not_found", "keyword": keyword}
        first_slide, idx = poses[0]
        snippet = _pptx_get_following_lines(slides, first_slide, idx, max_lines=context_lines)
        return {"status": "success", "keyword": keyword, "first_slide": first_slide, "snippet": snippet}
    except Exception as e:
        return {"error": "查找失败", "raw": f"{type(e).__name__}: {e}"}


@mcp.tool(description="仅分析 PDF 指定页集合，减少无关页干扰（页码从 1 开始）")
async def doc_analyze_pages(
    path: str = Field(description="本地 PDF 文件路径"),
    pages: List[int] = Field(description="要分析的页码列表（1-based）"),
    prompt: Optional[str] = Field(default=None, description="给 LLM 的额外提示")
):
    if not os.path.exists(path):
        return {"error": f"文件不存在: {path}"}
    ext = os.path.splitext(path)[-1].lower()
    if ext != ".pdf":
        return {"error": "doc_analyze_pages 仅支持 PDF"}
    try:
        doc = fitz.open(path)
        max_p = doc.page_count
        want = sorted({p for p in pages if isinstance(p, int) and 1 <= p <= max_p})
        if not want:
            return {"error": f"页码无效，合法范围 1..{max_p}"}
        content: List[dict] = []
        for p in want:
            page = doc[p - 1]
            pix = page.get_pixmap(dpi=_PDF_RENDER_DPI, alpha=False)
            png_bytes = _resize_image_bytes(pix.tobytes("png"), max_width=_DOC_MAX_WIDTH)
            b64 = base64.b64encode(png_bytes).decode()
            content.append({"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}"}})
        doc.close()
        content.append({"type": "text", "text": (_DEEP_ANALYZER_INSTRUCTION + (prompt or "请只输出答案，禁止解释。"))})
        messages = [{"role": "user", "content": content}]
        analysis = _call_vlm(messages)
        return {"status": "success", "path": path, "pages": want, "analysis": analysis}
    except Exception as e:
        return {"error": "分析失败", "raw": f"{type(e).__name__}: {e}"}


@mcp.tool(description="裁剪并分析 PDF/图片的矩形区域（用于聚焦表格/小字等）")
async def doc_analyze_region(
    path: str = Field(description="本地 PDF 或图片路径"),
    x: int = Field(description="裁剪矩形左上角 x 像素"),
    y: int = Field(description="裁剪矩形左上角 y 像素"),
    width: int = Field(description="裁剪矩形宽度像素"),
    height: int = Field(description="裁剪矩形高度像素"),
    page: int = Field(default=1, description="PDF 页码（1-based），图片忽略此参数"),
    prompt: Optional[str] = Field(default=None, description="给 LLM 的额外提示")
):
    if not os.path.exists(path):
        return {"error": f"文件不存在: {path}"}
    ext = os.path.splitext(path)[-1].lower()
    try:
        if ext == ".pdf":
            doc = fitz.open(path)
            if not (1 <= page <= doc.page_count):
                return {"error": f"页码无效，合法范围 1..{doc.page_count}"}
            pix = doc[page - 1].get_pixmap(dpi=_PDF_RENDER_DPI, alpha=False)
            img = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")
        elif ext in {".jpg", ".jpeg", ".png", ".bmp"}:
            img = Image.open(path).convert("RGB")
        else:
            return {"error": "仅支持 pdf / jpg / jpeg / png / bmp"}

        # 裁剪区域
        crop_box = (max(0, x), max(0, y), max(0, x + width), max(0, y + height))
        img_cropped = img.crop(crop_box)
        buf = io.BytesIO()
        img_cropped.save(buf, format="PNG", optimize=True)
        png_bytes = _resize_image_bytes(buf.getvalue(), max_width=_DOC_MAX_WIDTH)
        b64 = base64.b64encode(png_bytes).decode()
        content = [
            {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}"}},
            {"type": "text", "text": (_DEEP_ANALYZER_INSTRUCTION + (prompt or "请只输出答案，禁止解释。"))},
        ]
        messages = [{"role": "user", "content": content}]
        analysis = _call_vlm(messages)
        payload = {"status": "success", "path": path, "region": [x, y, width, height], "analysis": analysis}
        if ext == ".pdf":
            payload["page"] = page
        return payload
    except Exception as e:
        return {"error": "分析失败", "raw": f"{type(e).__name__}: {e}"}


@mcp.tool(description="获取文件信息：类型、页数/尺寸，辅助后续精细调用")
async def doc_get_info(path: str = Field(description="本地文件路径")):
    if not os.path.exists(path):
        return {"error": f"文件不存在: {path}"}
    ext = os.path.splitext(path)[-1].lower()
    info = {"path": path, "ext": ext}
    try:
        if ext == ".pdf":
            doc = fitz.open(path)
            info["type"] = "pdf"
            info["pages"] = doc.page_count
            doc.close()
        elif ext == ".pptx":
            info["type"] = "pptx"
            if Presentation is not None:
                prs = Presentation(path)
                info["slides"] = len(list(prs.slides))
        elif ext in {".jpg", ".jpeg", ".png", ".bmp"}:
            with Image.open(path) as im:
                w, h = im.size
            info["type"] = "image"
            info["width"], info["height"] = w, h
        else:
            info["type"] = "unknown"
        return {"status": "success", "info": info}
    except Exception as e:
        return {"error": "获取信息失败", "raw": f"{type(e).__name__}: {e}"}

# ======================= 启动 =======================
if __name__ == "__main__":
    mcp.run()
